"""
Generate CEO-facing documents:
1. PowerPoint overview presentation
2. Word document with detailed flow and usage
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGB, Inches as DocInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ── Colour palette ────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # electric blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
SUBTEXT   = RGBColor(0xA0, 0xB4, 0xC8)
GREEN     = RGBColor(0x00, 0xD4, 0x8A)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

PPT_PATH  = os.path.join(OUTPUT_DIR, "Agentic_AI_Website_Builder_Overview.pptx")
DOCX_PATH = os.path.join(OUTPUT_DIR, "Agentic_AI_Website_Builder_Detailed_Guide.docx")


# ═══════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════

def add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def accent_bar(slide, t, w=9.5, h=0.04):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.25), Inches(t), Inches(w), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def bullet_box(slide, items, l, t, w, h, size=14, color=WHITE, indent=0.25):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  ›  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def slide_num(prs):
    return len(prs.slides)


# ═══════════════════════════════════════════════════════════════════
#  POWERPOINT
# ═══════════════════════════════════════════════════════════════════

def build_ppt():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank

    # ── Slide 1 · Title ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    # gradient strip top
    accent_bar(s, 0.0, w=13.33, h=0.5)
    accent_bar(s, 7.0, w=13.33, h=0.5)

    txb(s, "AGENTIC AI", 0.5, 1.6, 12, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, "Website Builder Platform", 0.5, 2.8, 12, 0.8, size=30, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    accent_bar(s, 3.8, w=5, h=0.05)
    txb(s, "Executive Overview", 0.5, 4.0, 12, 0.6, size=20, color=SUBTEXT, align=PP_ALIGN.CENTER)
    txb(s, "Confidential  ·  May 2026", 0.5, 4.7, 12, 0.5, size=14, color=SUBTEXT, align=PP_ALIGN.CENTER, italic=True)

    # ── Slide 2 · Executive Summary ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Executive Summary", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    summary = (
        "Agentic AI Website Builder is an enterprise-grade SaaS platform that enables "
        "any organisation to generate, customise, and publish professional websites "
        "through a simple conversation — no developers required.\n\n"
        "Powered by a five-stage AI agent pipeline (CrewAI + OpenAI GPT-4.1-mini), "
        "the platform handles everything from requirements gathering to live hosting on AWS S3, "
        "with built-in authentication, Stripe billing, multi-tenant client management, "
        "and real-time monitoring — all in one product."
    )
    txb(s, summary, 0.4, 0.9, 12.5, 3.0, size=15, color=WHITE)

    # KPI boxes
    kpis = [
        ("5", "AI Agents"),
        ("< 2 min", "Website Generated"),
        ("∞", "Clients / Tenants"),
        ("AWS S3", "Instant Hosting"),
    ]
    for idx, (val, label) in enumerate(kpis):
        x = 0.5 + idx * 3.1
        box = s.shapes.add_shape(1, Inches(x), Inches(4.2), Inches(2.7), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        box.line.color.rgb = ACCENT; box.line.width = Pt(1.2)
        txb(s, val,   x, 4.35, 2.7, 0.8, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txb(s, label, x, 5.15, 2.7, 0.7, size=13, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # ── Slide 3 · The Problem ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "The Problem We Solve", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    problems = [
        "Building a professional website takes weeks and costs thousands of dollars",
        "SMBs and enterprises lack in-house web development talent",
        "Traditional website builders require technical expertise and design skills",
        "Keeping websites updated demands ongoing developer involvement",
        "Multi-tenant web agencies struggle to scale client site production",
    ]
    for i, prob in enumerate(problems):
        y = 1.1 + i * 1.1
        box = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        box.line.color.rgb = RGBColor(0xFF, 0x4D, 0x4D); box.line.width = Pt(0.8)
        txb(s, f"✕  {prob}", 0.65, y + 0.1, 12, 0.65, size=14, color=WHITE)

    # ── Slide 4 · Our Solution ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Our Solution", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    solutions = [
        ("Conversational AI", "Describe your website in plain English — agents handle the rest"),
        ("5-Stage Agent Pipeline", "Requirements → Design → Theme → Content → Code, fully automated"),
        ("One-Click Hosting", "Publish to AWS S3 in seconds; staging and live environments included"),
        ("Multi-Tenant SaaS", "Manage unlimited clients, each with isolated sites and billing"),
        ("Enterprise Security", "JWT auth, OTP (email + SMS), encrypted secret store, role-based access"),
    ]
    for i, (title, desc) in enumerate(solutions):
        y = 1.0 + i * 1.15
        box = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.3), Inches(0.9))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x1A)
        box.line.color.rgb = GREEN; box.line.width = Pt(0.8)
        txb(s, f"✓  {title}", 0.65, y + 0.05, 3.5, 0.5, size=14, bold=True, color=GREEN)
        txb(s, desc, 4.0, y + 0.08, 8.7, 0.65, size=13, color=WHITE)

    # ── Slide 5 · AI Agent Pipeline ───────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "The AI Agent Pipeline", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    agents = [
        ("1", "Requirements\nAnalyst", "Parses intent, industry,\npages & features"),
        ("2", "Designer\nAgent",       "Layout, colour palette\n& component planning"),
        ("3", "Theme\nAgent",          "CSS design system,\ntypography & branding"),
        ("4", "Content\nAgent",        "AI copywriting enhanced\nby live web search"),
        ("5", "Developer\nAgent",      "Generates production-\nready HTML/CSS/JS"),
    ]
    for i, (num, name, desc) in enumerate(agents):
        x = 0.35 + i * 2.55
        # circle badge
        circ = s.shapes.add_shape(9, Inches(x + 0.65), Inches(1.1), Inches(0.9), Inches(0.9))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        txb(s, num, x + 0.65, 1.12, 0.9, 0.7, size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        # card
        card = s.shapes.add_shape(1, Inches(x), Inches(2.1), Inches(2.3), Inches(3.2))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        card.line.color.rgb = ACCENT; card.line.width = Pt(1)
        txb(s, name, x, 2.2, 2.3, 0.9, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txb(s, desc, x, 3.1, 2.3, 1.5, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        # arrow
        if i < 4:
            arrow = s.shapes.add_shape(1, Inches(x + 2.32), Inches(2.55), Inches(0.22), Inches(0.4))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    txb(s, "↓  Output: Hosted on AWS S3 within minutes", 0.5, 5.5, 12.3, 0.6,
        size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # ── Slide 6 · Platform Architecture ──────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Platform Architecture", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    layers = [
        ("Frontend", "Login · Dashboard · Console · Monitoring · Logs  (HTML/JS)"),
        ("API Layer", "FastAPI REST — Auth · Website Builder · Payments · Commerce · Monitoring"),
        ("Agent Engine", "CrewAI — 5 specialised agents orchestrated per request"),
        ("Services", "Auth · Hosting · Payment · OTP · Notifications · Analytics · Currency"),
        ("Data & Storage", "Snowflake (prod) / SQLite (dev) · AWS S3 · Encrypted Secret Store"),
    ]
    colors = [RGBColor(0x1A,0x3A,0x5C), RGBColor(0x14,0x30,0x48), RGBColor(0x0E,0x26,0x3C),
              RGBColor(0x0A,0x1E,0x30), RGBColor(0x06,0x16,0x24)]
    for i, (layer, detail) in enumerate(layers):
        y = 1.0 + i * 1.1
        box = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.9))
        box.fill.solid(); box.fill.fore_color.rgb = colors[i]
        box.line.color.rgb = ACCENT; box.line.width = Pt(0.6)
        txb(s, layer,  0.55, y + 0.1, 2.0, 0.65, size=13, bold=True, color=ACCENT)
        txb(s, detail, 2.6,  y + 0.1, 10.0, 0.65, size=12, color=WHITE)

    # ── Slide 7 · Key Features ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Key Platform Features", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    cols = [
        ("🔐  Security & Auth",
         ["JWT + OTP (Email & SMS)", "Bcrypt password hashing", "Role-based access control",
          "Encrypted secret store", "Rate limiting (SlowAPI)"]),
        ("💳  Billing & Commerce",
         ["Stripe subscriptions", "Pro & Enterprise plans", "Shopping cart module",
          "Catalogue scraper", "Payment reminders"]),
        ("☁️  Hosting & DevOps",
         ["AWS S3 staging + publish", "Custom domain support", "Scheduled monitoring",
          "Real-time health dashboard", "Uptime alerting"]),
        ("🤖  AI Capabilities",
         ["GPT-4.1-mini powered", "Web search (Tavily)", "Social search (Twitter)",
          "Website scraping for context", "Multi-language content"]),
    ]
    for i, (title, items) in enumerate(cols):
        x = 0.35 + i * 3.2
        card = s.shapes.add_shape(1, Inches(x), Inches(1.05), Inches(3.0), Inches(5.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        card.line.color.rgb = ACCENT; card.line.width = Pt(1)
        txb(s, title, x, 1.1, 3.0, 0.65, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            txb(s, f"›  {item}", x + 0.1, 1.85 + j * 0.85, 2.8, 0.7, size=12, color=WHITE)

    # ── Slide 8 · Use Cases ───────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Who Benefits", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    cases = [
        ("🏥 Healthcare", "Clinics, diagnostics labs, and hospitals launch patient-facing sites in minutes"),
        ("🏪 Retail & D2C", "Brands go live with product catalogues and shopping carts instantly"),
        ("🎓 Education", "Schools and tutors publish course pages with structured content"),
        ("⚖️ Professional Services", "Law firms, consultants, and HR agencies get polished landing pages"),
        ("🏨 Hospitality", "Hotels, salons, and restaurants showcase services with AI-curated copy"),
        ("🏢 Enterprises", "IT teams spin up microsites and campaign pages without developer queues"),
    ]
    for i, (sector, desc) in enumerate(cases):
        row, col = divmod(i, 3)
        x = 0.4 + col * 4.3
        y = 1.05 + row * 2.7
        card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.0), Inches(2.3))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        card.line.color.rgb = ACCENT; card.line.width = Pt(0.8)
        txb(s, sector, x, y + 0.1, 4.0, 0.6, size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txb(s, desc,   x + 0.15, y + 0.75, 3.7, 1.3, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 9 · Competitive Advantage ──────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Why We Win", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    headers = ["", "Agentic AI Builder", "Wix / Squarespace", "Freelance Agency", "In-House Dev"]
    rows = [
        ["AI-generated content",    "✓", "✗", "~", "~"],
        ["Multi-agent pipeline",    "✓", "✗", "✗", "✗"],
        ["Live web search context", "✓", "✗", "~", "✗"],
        ["Built-in multi-tenancy",  "✓", "✗", "✗", "~"],
        ["Stripe + OTP built-in",   "✓", "✗", "✗", "~"],
        ["Minutes to publish",      "✓", "~", "✗", "✗"],
        ["Self-hosted / private",   "✓", "✗", "✗", "✓"],
    ]
    col_w = [3.0, 2.3, 2.3, 2.3, 2.3]
    col_x = [0.35, 3.4, 5.75, 8.1, 10.45]
    # header row
    for ci, (hdr, cx) in enumerate(zip(headers, col_x)):
        hbox = s.shapes.add_shape(1, Inches(cx), Inches(1.05), Inches(col_w[ci]), Inches(0.55))
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = ACCENT if ci == 1 else RGBColor(0x1A, 0x3A, 0x5C)
        hbox.line.color.rgb = DARK_BG; hbox.line.width = Pt(0.5)
        txb(s, hdr, cx, 1.1, col_w[ci], 0.45, size=12, bold=True,
            color=DARK_BG if ci == 1 else WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        for ci, (cell, cx) in enumerate(zip(row, col_x)):
            y = 1.65 + ri * 0.73
            rbox = s.shapes.add_shape(1, Inches(cx), Inches(y), Inches(col_w[ci]), Inches(0.63))
            rbox.fill.solid()
            rbox.fill.fore_color.rgb = (RGBColor(0x0A,0x2A,0x1A) if ci == 1
                                        else RGBColor(0x10,0x24,0x36))
            rbox.line.color.rgb = RGBColor(0x1E,0x3A,0x54); rbox.line.width = Pt(0.4)
            cell_color = GREEN if cell == "✓" else (RGBColor(0xFF,0x6B,0x6B) if cell == "✗" else SUBTEXT)
            txb(s, cell, cx, y + 0.05, col_w[ci], 0.5, size=13,
                bold=(ci == 1), color=cell_color if ci != 0 else WHITE,
                align=PP_ALIGN.CENTER if ci != 0 else PP_ALIGN.LEFT)

    # ── Slide 10 · Deployment / Tech Stack ───────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.72)
    txb(s, "Technology Stack", 0.4, 0.15, 12, 0.55, size=26, bold=True, color=ACCENT)

    stack = [
        ("AI / NLP",      "OpenAI GPT-4.1-mini · CrewAI 1.14 · Tavily Web Search"),
        ("Backend",       "Python 3.9+ · FastAPI 0.115 · Uvicorn · Pydantic v2"),
        ("Database",      "Snowflake (production) · SQLite (development fallback)"),
        ("Cloud",         "AWS S3 (hosting) · Google Drive (optional image store)"),
        ("Auth & Security","JWT · bcrypt · Fernet encryption · SlowAPI rate limiting"),
        ("Payments",      "Stripe 12.1 · Webhook verification · Subscription management"),
        ("Notifications", "Twilio SMS & WhatsApp · SendGrid (SMTP) Email OTP"),
        ("Frontend",      "Vanilla HTML/CSS/JS · No frontend build step required"),
    ]
    for i, (cat, tech) in enumerate(stack):
        row, col = divmod(i, 2)
        x = 0.4 + col * 6.5
        y = 1.05 + row * 1.45
        card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(6.1), Inches(1.2))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E)
        card.line.color.rgb = ACCENT; card.line.width = Pt(0.8)
        txb(s, cat,  x + 0.15, y + 0.05, 2.0,  0.5, size=13, bold=True, color=ACCENT)
        txb(s, tech, x + 0.15, y + 0.55, 5.7, 0.55, size=12, color=WHITE)

    # ── Slide 11 · Call to Action ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK_BG)
    accent_bar(s, 0.0, w=13.33, h=0.5)
    accent_bar(s, 7.0, w=13.33, h=0.5)

    txb(s, "Ready to Transform", 0.5, 1.3, 12.3, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, "Your Digital Presence?", 0.5, 2.3, 12.3, 1.0, size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s,
        "From conversation to live website in under 2 minutes.\n"
        "Enterprise-grade security, billing, and multi-tenant management — out of the box.",
        0.5, 3.5, 12.3, 1.2, size=16, color=SUBTEXT, align=PP_ALIGN.CENTER)

    actions = ["Schedule a Live Demo", "Request Pilot Access", "Review Technical Documentation"]
    for i, action in enumerate(actions):
        x = 1.0 + i * 3.9
        btn = s.shapes.add_shape(1, Inches(x), Inches(5.0), Inches(3.4), Inches(0.75))
        btn.fill.solid(); btn.fill.fore_color.rgb = ACCENT if i == 0 else RGBColor(0x14, 0x2A, 0x3E)
        btn.line.color.rgb = ACCENT; btn.line.width = Pt(1.2)
        txb(s, action, x, 5.08, 3.4, 0.55, size=13, bold=(i==0),
            color=DARK_BG if i == 0 else WHITE, align=PP_ALIGN.CENTER)

    txb(s, "senthilvasansubbu  ·  agentic-s3-website-builder  ·  Confidential",
        0.5, 6.2, 12.3, 0.4, size=11, color=SUBTEXT, align=PP_ALIGN.CENTER, italic=True)

    prs.save(PPT_PATH)
    print(f"✅  PPT saved → {PPT_PATH}")


# ═══════════════════════════════════════════════════════════════════
#  WORD DOCUMENT
# ═══════════════════════════════════════════════════════════════════

def set_doc_style(doc):
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = DocPt(11)


def h1(doc, text):
    p = doc.add_heading(text, level=1)
    p.runs[0].font.color.rgb = DocRGB(0x00, 0x5A, 0x9E)
    return p


def h2(doc, text):
    p = doc.add_heading(text, level=2)
    p.runs[0].font.color.rgb = DocRGB(0x00, 0x7A, 0xCC)
    return p


def h3(doc, text):
    p = doc.add_heading(text, level=3)
    p.runs[0].font.color.rgb = DocRGB(0x00, 0x99, 0xDD)
    return p


def para(doc, text, bold=False, italic=False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    p.paragraph_format.space_after = DocPt(6)
    return p


def bullet(doc, items, level=0):
    for item in items:
        p = doc.add_paragraph(item, style='List Bullet')
        p.paragraph_format.space_after = DocPt(4)
        p.paragraph_format.left_indent = DocInches(0.25 + level * 0.25)


def add_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        run = hdr_cells[i].paragraphs[0].runs[0]
        run.bold = True
        hdr_cells[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        tc = hdr_cells[i]._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), '00589E')
        tcPr.append(shd)
        run.font.color.rgb = DocRGB(0xFF, 0xFF, 0xFF)
    for ri, row in enumerate(rows):
        cells = table.rows[ri + 1].cells
        for ci, cell in enumerate(row):
            cells[ci].text = str(cell)
            cells[ci].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT
    doc.add_paragraph()


def build_docx():
    doc = Document()
    set_doc_style(doc)

    # ── Cover ───────────────────────────────────────────────────────
    doc.add_paragraph()
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run("Agentic AI Website Builder Platform")
    title_run.font.size = DocPt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = DocRGB(0x00, 0x5A, 0x9E)

    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub_p.add_run("Detailed Product Guide — Flow, Usage & Architecture")
    sub_run.font.size = DocPt(16)
    sub_run.italic = True
    sub_run.font.color.rgb = DocRGB(0x44, 0x72, 0xA4)

    date_p = doc.add_paragraph()
    date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_p.add_run("Confidential  ·  May 2026").font.color.rgb = DocRGB(0x80, 0x80, 0x80)

    doc.add_paragraph()
    doc.add_page_break()

    # ── 1. Introduction ──────────────────────────────────────────────
    h1(doc, "1. Introduction")
    para(doc,
         "Agentic AI Website Builder is a fully managed, multi-tenant SaaS platform that empowers "
         "organisations to design, generate, and publish production-ready websites through natural "
         "language conversation. The platform eliminates the need for in-house web developers, "
         "graphic designers, or copywriters by orchestrating five specialised AI agents — each "
         "responsible for a distinct stage of the website creation pipeline.")
    para(doc,
         "Built on FastAPI and CrewAI (OpenAI GPT-4.1-mini), the system is production-ready with "
         "enterprise features including JWT authentication, OTP-based login, Stripe billing, "
         "AWS S3 hosting, real-time monitoring, and encrypted secret management.")

    h2(doc, "1.1 Target Audience")
    bullet(doc, [
        "CEOs and CXOs evaluating AI-powered digital transformation tools",
        "IT/Tech leads assessing platform architecture and integration points",
        "Product and growth teams exploring rapid website deployment capabilities",
        "Agency owners looking for a white-label, multi-tenant website builder",
    ])

    # ── 2. Platform Overview ──────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "2. Platform Overview")
    para(doc,
         "The platform operates as a RESTful API backend (FastAPI + Uvicorn) with a lightweight "
         "HTML/JS frontend. Clients interact via a browser-based dashboard, while the AI agents "
         "operate server-side, processing requests asynchronously.")

    h2(doc, "2.1 High-Level Architecture")
    layers = [
        ("Presentation Layer", "login.html, dashboard.html, console.html, monitoring.html, logs.html"),
        ("API Layer", "FastAPI routes — auth, website_builder, chatbot, clients, payment, commerce, monitoring, admin"),
        ("Agent Engine", "CrewAI orchestrator running 5 sequential AI agents"),
        ("Service Layer", "Auth, Hosting, Payment, OTP, Notifications, Analytics, Currency, Secret Store"),
        ("Data Layer", "Snowflake (production) / SQLite (development) + AWS S3 file storage"),
    ]
    add_table(doc, ["Layer", "Components"], layers, [2.5, 4.5])

    h2(doc, "2.2 Core Technology Stack")
    tech_rows = [
        ("AI / NLP",       "OpenAI GPT-4.1-mini, CrewAI 1.14, Tavily Web Search"),
        ("Backend",        "Python 3.9+, FastAPI 0.115, Uvicorn 0.34, Pydantic v2"),
        ("Database",       "Snowflake (prod), SQLite (dev auto-fallback)"),
        ("Cloud / Storage","AWS S3, Google Drive (optional image backend)"),
        ("Auth / Security","JWT (PyJWT), bcrypt, Fernet encryption, SlowAPI rate limiting"),
        ("Payments",       "Stripe 12.1 — subscriptions, webhooks, reminders"),
        ("Notifications",  "Twilio (SMS + WhatsApp), SendGrid / SMTP (Email OTP)"),
        ("Frontend",       "Vanilla HTML/CSS/JS — zero build toolchain"),
    ]
    add_table(doc, ["Category", "Technologies"], tech_rows)

    # ── 3. User Journey ───────────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "3. End-to-End User Journey")

    steps = [
        ("Step 1 — Registration & Login",
         "A new client signs up via the dashboard. The system sends an OTP to their email or "
         "mobile number (via SendGrid or Twilio). After OTP verification, a JWT token is issued "
         "and the client is redirected to the main dashboard."),
        ("Step 2 — Describe Your Website",
         "The client opens the chatbot interface and describes the website they need in plain English. "
         "Example: 'Build a diagnostic lab website for my clinic in Mumbai with appointment booking, "
         "service list, and contact form.' The system accepts references URLs, logo uploads, and "
         "preferred colour schemes as additional inputs."),
        ("Step 3 — Agent Pipeline Executes",
         "The Requirements Analyst agent parses the prompt and extracts structured requirements. "
         "The Designer Agent produces a layout and visual specification. The Theme Agent generates "
         "a CSS design system. The Content Agent writes all copy, enhanced by live Tavily web search. "
         "The Developer Agent assembles the final HTML/CSS/JS file."),
        ("Step 4 — Preview in Staging",
         "The generated site is saved to output/staging/<slug>/ and served as a preview URL. "
         "The client can review the site, request changes via the chatbot, and iterate until satisfied."),
        ("Step 5 — Publish to Production",
         "A single click publishes the site to AWS S3. The platform assigns a public URL, "
         "creates an immutable published snapshot in output/published/, and logs the deployment."),
        ("Step 6 — Ongoing Management",
         "The client can return to the dashboard at any time to update content, change themes, "
         "add pages, or re-publish. The monitoring service checks site uptime every 5 minutes "
         "and alerts via the dashboard or WhatsApp if a site goes down."),
    ]
    for title, desc in steps:
        h2(doc, title)
        para(doc, desc)

    # ── 4. AI Agent Pipeline ──────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "4. AI Agent Pipeline — Deep Dive")
    para(doc,
         "Each agent in the pipeline is a specialised CrewAI agent with a defined role, goal, "
         "backstory, and set of tools. Agents pass structured outputs to the next agent, ensuring "
         "coherent, high-quality results at every stage.")

    agents = [
        ("4.1 Requirements Analyst",
         "agents/requirements_analyst.py",
         "Parses the client's free-text prompt and extracts structured requirements: business type, "
         "target audience, number of pages, required features (booking, e-commerce, contact form, etc.), "
         "tone of voice, and brand preferences. Outputs a structured JSON requirements document.",
         ["Accepts: raw user prompt + optional reference URL",
          "Outputs: structured requirements object",
          "Tools: website_scraper (for reference URLs), web_search (Tavily)"]),

        ("4.2 Designer Agent",
         "agents/designer_agent.py",
         "Takes the requirements object and produces a detailed design specification: page layout "
         "for each section, colour palette (primary, secondary, accent, background, text), "
         "typography (font families, sizes, weights), and a component list (hero, nav, cards, footer).",
         ["Accepts: requirements object",
          "Outputs: design specification document",
          "Tools: social_media_search (for design inspiration trends)"]),

        ("4.3 Theme Agent",
         "agents/theme_agent.py",
         "Converts the design specification into a concrete CSS design system: CSS custom properties "
         "(variables), utility classes, responsive breakpoints, and animation definitions. "
         "The theme is injected into the final HTML as an embedded stylesheet.",
         ["Accepts: design specification",
          "Outputs: CSS theme string",
          "Tools: theme_builder"]),

        ("4.4 Content Agent",
         "agents/content_agent.py",
         "Writes all website copy — headlines, section text, CTAs, team bios, service descriptions, "
         "FAQs, and footer content — tailored to the business context. Uses Tavily web search to "
         "gather real industry data and social search to align with current trends.",
         ["Accepts: requirements + design spec",
          "Outputs: structured content map (section → copy)",
          "Tools: web_search, social_media_search, website_scraper"]),

        ("4.5 Developer Agent",
         "agents/developer_agent.py",
         "Assembles the complete, production-ready HTML file by combining the CSS theme, structured "
         "content, and component templates. Generates semantic HTML5, responsive CSS3, and vanilla "
         "JavaScript for interactive elements (sliders, modals, form validation, smooth scroll).",
         ["Accepts: theme + content map + design spec",
          "Outputs: single self-contained HTML/CSS/JS file",
          "Tools: html_generator, s3_uploader"]),
    ]
    for section, file, desc, bullets in agents:
        h2(doc, section)
        para(doc, f"File: {file}", italic=True)
        para(doc, desc)
        bullet(doc, bullets)
        doc.add_paragraph()

    # ── 5. API Reference ──────────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "5. API Reference")
    para(doc, "Full interactive documentation is available at http://localhost:8000/docs")

    api_rows = [
        ("POST", "/auth/register",              "Register a new client account"),
        ("POST", "/auth/login",                 "Authenticate; returns JWT access token"),
        ("POST", "/auth/otp/send",              "Send OTP to email or SMS"),
        ("POST", "/auth/otp/verify",            "Verify OTP and activate session"),
        ("POST", "/website-builder/generate",   "Generate website from prompt (triggers agent pipeline)"),
        ("GET",  "/website-builder/sites",      "List all websites for current client"),
        ("GET",  "/website-builder/sites/{id}", "Get details of a specific site"),
        ("POST", "/website-builder/publish",    "Publish staging site to AWS S3"),
        ("POST", "/chatbot/message",            "Send a chat message to iterate on the website"),
        ("GET",  "/monitoring/status",          "Get platform and site uptime status"),
        ("POST", "/payment/checkout",           "Create Stripe checkout session"),
        ("POST", "/payment/webhook",            "Handle Stripe webhook events"),
        ("GET",  "/admin/clients",              "List all clients (superuser only)"),
        ("POST", "/commerce/cart",              "Add product to shopping cart"),
        ("GET",  "/commerce/catalogue",         "List product catalogue for a client site"),
    ]
    add_table(doc, ["Method", "Endpoint", "Description"], api_rows)

    # ── 6. Security Architecture ──────────────────────────────────────
    doc.add_page_break()
    h1(doc, "6. Security Architecture")

    h2(doc, "6.1 Authentication Flow")
    para(doc,
         "1.  Client submits email + password to POST /auth/login.\n"
         "2.  Server verifies password hash (bcrypt) and issues a signed JWT (HS256).\n"
         "3.  JWT is stored client-side and sent in the Authorization: Bearer header.\n"
         "4.  Token expiry is configurable via JWT_TTL_MINUTES (default: 1440 minutes / 24 hours).\n"
         "5.  OTP is required for first login and password reset — delivered via SendGrid or Twilio.")

    h2(doc, "6.2 Secret Store")
    para(doc,
         "Per-client API keys (e.g., Stripe keys, custom AWS credentials) are encrypted at rest "
         "using Fernet symmetric encryption (AES-128-CBC + HMAC). The master key is set via "
         "STORAGE_SECRETS_KEY in the environment and never stored in the database.")

    h2(doc, "6.3 Rate Limiting")
    para(doc,
         "SlowAPI middleware enforces per-IP rate limits on all public endpoints to prevent "
         "brute-force and denial-of-service attacks.")

    h2(doc, "6.4 Data Isolation")
    para(doc,
         "All database queries are scoped by client_id. Row-level isolation ensures that no client "
         "can access another's websites, credentials, or billing data.")

    # ── 7. Deployment Guide ───────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "7. Deployment Guide")

    h2(doc, "7.1 Local / Development Setup")
    para(doc, "Prerequisites: Python 3.9+, Git, AWS credentials (optional)")
    bullet(doc, [
        "git clone https://github.com/senthilvasansubbu/agentic-s3-website-builder.git",
        "cd agentic-s3-website-builder",
        "pip install -r requirements.txt",
        "cp .env.example .env  # Fill in OPENAI_API_KEY and JWT_SECRET at minimum",
        "python app.py         # Server starts at http://localhost:8000",
    ])

    h2(doc, "7.2 Environment Variables (Key)")
    env_rows = [
        ("OPENAI_API_KEY",         "Required", "OpenAI API key for GPT-4.1-mini"),
        ("JWT_SECRET",             "Required", "Long random string for JWT signing"),
        ("STORAGE_SECRETS_KEY",    "Required", "Fernet key for secret store encryption"),
        ("AWS_ACCESS_KEY_ID",      "Optional", "AWS credentials for S3 hosting"),
        ("SNOWFLAKE_ACCOUNT",      "Optional", "Snowflake for production database"),
        ("STRIPE_SECRET_KEY",      "Optional", "Stripe for billing features"),
        ("SMTP_PASSWORD",          "Optional", "SendGrid / SMTP for email OTP"),
        ("TWILIO_AUTH_TOKEN",      "Optional", "Twilio for SMS / WhatsApp OTP"),
        ("TAVILY_API_KEY",         "Optional", "Tavily for live web search in content agent"),
        ("MONITOR_INTERVAL_MINUTES","Optional","Uptime check frequency (default: 5)"),
    ]
    add_table(doc, ["Variable", "Required?", "Purpose"], env_rows)

    h2(doc, "7.3 Production Recommendations")
    bullet(doc, [
        "Run behind a reverse proxy (Nginx / Caddy) with TLS termination",
        "Set PLATFORM_API_URL to your public domain for external monitoring checks",
        "Use Snowflake for persistent, scalable production data storage",
        "Enable AWS S3 bucket versioning for published site recovery",
        "Rotate JWT_SECRET and STORAGE_SECRETS_KEY periodically",
        "Restrict admin endpoints to IP allowlist at the proxy level",
    ])

    # ── 8. Monitoring & Analytics ──────────────────────────────────────
    doc.add_page_break()
    h1(doc, "8. Monitoring & Analytics")
    para(doc,
         "The built-in monitoring service (APScheduler) runs configurable uptime checks against "
         "all published client websites and the platform API itself. Results are stored and surfaced "
         "via the monitoring dashboard (monitoring.html) and the GET /monitoring/status endpoint.")

    h2(doc, "8.1 Monitoring Configuration")
    bullet(doc, [
        "MONITOR_INTERVAL_MINUTES — how often checks run (default: every 5 minutes)",
        "MONITOR_TIMEOUT_SEC — HTTP request timeout per check (default: 10 seconds)",
        "PLATFORM_API_URL — the platform's own public URL, used for self-health checks",
    ])

    h2(doc, "8.2 Analytics Service")
    para(doc,
         "The analytics_service.py module tracks key usage metrics: websites generated per client, "
         "publish events, login activity, and chatbot session volume. Data is queryable via the "
         "admin console for business intelligence reporting.")

    # ── 9. Commerce & Billing ──────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "9. Commerce & Billing")

    h2(doc, "9.1 Stripe Integration")
    para(doc,
         "Platform subscriptions (Pro and Enterprise plans) are managed through Stripe. "
         "Clients select a plan, are redirected to a Stripe Checkout session, and upon successful "
         "payment are upgraded in the database. Webhooks handle subscription renewals, failures, "
         "and cancellations automatically.")

    h2(doc, "9.2 Shopping Cart & Catalogue")
    para(doc,
         "For clients building e-commerce websites, the platform includes a shopping cart module "
         "and a catalogue scraper that can import products from existing websites. Cart state is "
         "persisted per session and can be embedded directly into generated websites.")

    h2(doc, "9.3 Payment Reminders")
    para(doc,
         "The payment_reminder_service.py sends automated email and WhatsApp reminders to clients "
         "with overdue invoices, reducing churn and manual follow-up overhead.")

    # ── 10. Roadmap ───────────────────────────────────────────────────
    doc.add_page_break()
    h1(doc, "10. Suggested Roadmap")

    roadmap = [
        ("Q3 2026", "Custom domain mapping (CNAME → S3/CloudFront)"),
        ("Q3 2026", "White-label reseller portal for agencies"),
        ("Q3 2026", "Image generation integration (DALL·E / Stable Diffusion)"),
        ("Q4 2026", "Mobile app companion (iOS / Android) for on-the-go edits"),
        ("Q4 2026", "A/B testing module for landing page variants"),
        ("Q1 2027", "Multi-language website generation (i18n)"),
        ("Q1 2027", "CRM integration (HubSpot, Salesforce) for lead capture"),
        ("Q2 2027", "Self-serve plan upgrades and usage-based billing"),
    ]
    add_table(doc, ["Timeline", "Feature"], roadmap)

    # ── 11. Glossary ──────────────────────────────────────────────────
    h1(doc, "11. Glossary")
    glossary = [
        ("Agent",      "An autonomous AI role within the CrewAI framework, with a specific responsibility"),
        ("CrewAI",     "Open-source framework for orchestrating multiple AI agents in sequential or parallel pipelines"),
        ("JWT",        "JSON Web Token — a signed, stateless authentication credential"),
        ("OTP",        "One-Time Password — a single-use code sent via email or SMS for verification"),
        ("S3",         "Amazon Simple Storage Service — used for hosting static website files"),
        ("Staging",    "A preview environment where generated sites are reviewed before publication"),
        ("Snowflake",  "Cloud data warehouse used as the production database"),
        ("Fernet",     "Symmetric authenticated encryption scheme used for the secret store"),
        ("Tavily",     "AI-powered web search API used by the Content Agent to gather real-world information"),
        ("Slug",       "URL-friendly identifier for a website (e.g., 'my-clinic-mumbai')"),
    ]
    add_table(doc, ["Term", "Definition"], glossary)

    doc.save(DOCX_PATH)
    print(f"✅  DOCX saved → {DOCX_PATH}")


if __name__ == "__main__":
    build_ppt()
    build_docx()
    print("\n🎉  Both documents generated successfully.")
